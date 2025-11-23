"""
Generate PowerPoint Presentation for scASDC Paper Analysis
Extracts objective functions, loss functions, and evaluation techniques
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_scASDC_presentation():
    """Create comprehensive PowerPoint presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    title_color = RGBColor(0, 51, 102)      # Dark blue
    accent_color = RGBColor(0, 112, 192)    # Blue
    highlight_color = RGBColor(237, 125, 49) # Orange
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "scASDC: Attention Enhanced Structural Deep Clustering"
    title.font.size = Pt(36)
    title.font.bold = True
    title.font.color.rgb = title_color
    title.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle = subtitle_frame.add_paragraph()
    subtitle.text = "Single-cell RNA-seq Data Analysis"
    subtitle.font.size = Pt(24)
    subtitle.font.color.rgb = accent_color
    subtitle.alignment = PP_ALIGN.CENTER
    
    # Add authors
    authors_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    authors_frame = authors_box.text_frame
    authors = authors_frame.add_paragraph()
    authors.text = "Wenwen Min, Zhen Wang, Fangfang Zhu, Taosheng Xu, Shunfang Wang"
    authors.font.size = Pt(14)
    authors.alignment = PP_ALIGN.CENTER
    
    # Add conference
    conf_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    conf_frame = conf_box.text_frame
    conf = conf_frame.add_paragraph()
    conf.text = "IEEE International Conference on Bioinformatics and Biomedicine (BIBM) 2024"
    conf.font.size = Pt(12)
    conf.font.italic = True
    conf.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Overview of scASDC Method"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Four Core Modules:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.space_after = Pt(10)
    
    modules = [
        "1. ZINB-based Autoencoder Module",
        "   • Learns latent representations of gene expression",
        "   • Models zero-inflation and overdispersion",
        "",
        "2. Graph Autoencoder Module (GCN)",
        "   • Captures high-order structural relationships",
        "   • Extracts cell-cell interaction patterns",
        "",
        "3. Attention Fusion Module",
        "   • Multi-head attention (8 heads)",
        "   • Layer-by-layer integration of content + structure",
        "",
        "4. Self-Supervised Learning Module",
        "   • End-to-end clustering training",
        "   • Dual KL divergence losses"
    ]
    
    for item in modules:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0 if item.startswith(('1', '2', '3', '4')) else 1
        p.font.size = Pt(16) if item.startswith(('1', '2', '3', '4')) else Pt(14)
        p.space_after = Pt(3)
    
    # Slide 3: Module 1 - ZINB Autoencoder
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Module 1: ZINB-based Autoencoder"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Objective:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "Learn low-dimensional representations capturing gene expression patterns"
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Mathematical Formulation:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    formulas = [
        "Encoder: H_l = ReLU(W_l × H_(l-1) + b_l)",
        "",
        "ZINB Parameters:",
        "  • Dropout: Π = sigmoid(W_π × H_L)",
        "  • Mean: M_i = diag(S_i) × exp(W_μ × H_L)",
        "  • Dispersion: Θ = exp(W_θ × H_L)",
        "",
        "Architecture: 2000 → 1000 → 1000 → 4000 → 10 nodes"
    ]
    
    for item in formulas:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.name = "Courier New" if "→" in item or "×" in item else "Calibri"
        p.space_after = Pt(2)
    
    # Slide 4: Module 2 - Graph Autoencoder
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Module 2: Graph Autoencoder (GCN)"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Objective:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "Capture high-order structural relationships between cells"
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Graph Convolutional Network Layer:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "Z_l = ReLU(D̂^(-1/2) × (A+I) × D̂^(-1/2) × R_(l-1) × U_(l-1))"
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Key Features:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    features = [
        "• A = KNN adjacency matrix (k=10)",
        "• I = identity matrix",
        "• D̂ = degree matrix",
        "• R_(l-1) = fused representation (attention output)",
        "• Mitigates oversmoothing problem in deep GCNs"
    ]
    
    for item in features:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(3)
    
    # Slide 5: Module 3 - Attention Fusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Module 3: Attention Fusion Mechanism"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Objective:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "Intelligently integrate gene expression + structural information"
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Multi-Head Attention (8 heads):"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    attention_items = [
        "Step 1: Weighted combination",
        "  Y_(l-1) = α × H_(l-1) + (1-α) × Z_(l-1)",
        "  where α = 0.5 (balance parameter)",
        "",
        "Step 2: Multi-head attention",
        "  head_i = softmax(Q × K^T / √d_k) × V",
        "  R_l = W × Concat(head_1, ..., head_8)",
        "",
        "Benefits:",
        "  • Adaptive feature weighting",
        "  • Layer-by-layer fusion prevents information loss",
        "  • Combines complementary information sources"
    ]
    
    for item in attention_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0 if item.startswith(('Step', 'Benefits:')) else 1
        p.font.size = Pt(14) if item.startswith(('Step', 'Benefits:')) else Pt(13)
        p.font.name = "Courier New" if "×" in item or "√" in item else "Calibri"
        p.space_after = Pt(2)
    
    # Slide 6: Module 4 - Self-Supervised Learning
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Module 4: Self-Supervised Learning"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Objective:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "Enable end-to-end clustering through soft assignments"
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Soft Assignment (Student's t-distribution):"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "q_ij = [1 + ||h_i - μ_j||²/λ]^(-(λ+1)/2) / Σ_j' [...]"
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Target Distribution (high-confidence):"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "p_ij = (q_ij² / g_j) / Σ_j' (q_ij'² / g_j')"
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "where g_j = Σ_i q_ij (soft cluster frequency)"
    p.font.size = Pt(12)
    p.font.italic = True
    
    # Slide 7: Loss Functions Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Five Loss Functions"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Total Loss Function:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "L_total = λ₁L_GAEX + λ₂L_GAEA + λ₃L_clu + λ₄L_gae + λ₅L_ZINB"
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.space_after = Pt(15)
    
    losses = [
        ("1. L_ZINB", "ZINB reconstruction loss", "λ₅ = 0.5"),
        ("2. L_GAEX", "Graph expression preservation", "λ₁ = 0.5"),
        ("3. L_GAEA", "Graph adjacency reconstruction", "λ₂ = 0.01"),
        ("4. L_clu", "Clustering loss (ZINB module)", "λ₃ = 0.1"),
        ("5. L_gae", "Clustering loss (GCN module)", "λ₄ = 0.01")
    ]
    
    for name, desc, weight in losses:
        p = tf.add_paragraph()
        p.text = f"{name}: {desc}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = f"   Weight: {weight}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.space_after = Pt(8)
    
    # Slide 8: Loss 1 - ZINB Loss
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Loss 1: ZINB Reconstruction Loss"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Purpose:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "Model the zero-inflated and overdispersed nature of scRNA-seq data"
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Formula:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "L_ZINB = -log(ZINB(X̄|π, μ, θ))"
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "ZINB Distribution:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "ZINB(X̄|π,μ,θ) = π·δ₀(X̄) + (1-π)·NB(X̄|μ,θ)"
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Key Properties:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    properties = [
        "• π: Dropout probability (models zero-inflation)",
        "• μ: Mean parameter (expression level)",
        "• θ: Dispersion parameter (overdispersion)",
        "• Negative log-likelihood guides distribution learning"
    ]
    
    for item in properties:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.space_after = Pt(3)
    
    # Slide 9: Loss 2 & 3 - Graph Losses
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Loss 2 & 3: Graph Autoencoder Losses"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Loss 2: Graph Adjacency Reconstruction (L_GAEA)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "L_GAEA = ||A - Â||²_F"
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.space_after = Pt(5)
    
    p = tf.add_paragraph()
    p.text = "where Â = sigmoid(Z_L^T × Z_L)"
    p.font.size = Pt(12)
    p.font.italic = True
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Purpose: Preserve structural relationships between cells"
    p.font.size = Pt(13)
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Loss 3: Graph Expression Preservation (L_GAEX)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "L_GAEX = ||X̄ - Z_L||²_F"
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Purpose: Ensure GCN output retains gene expression information"
    p.font.size = Pt(13)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Together: Structure + Content preservation"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = highlight_color
    
    # Slide 10: Loss 4 & 5 - Clustering Losses
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Loss 4 & 5: Clustering Losses (KL Divergence)"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Loss 4: ZINB Module Clustering (L_clu)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "L_clu = KL(P||Q) = Σᵢⱼ p_ij × log(p_ij / q_ij)"
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Applied to: H_(L/2) (middle layer of ZINB autoencoder)"
    p.font.size = Pt(12)
    p.font.italic = True
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Loss 5: GCN Module Clustering (L_gae)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "L_gae = KL(P||Z_pre) = Σᵢⱼ p_ij × log(p_ij / z_ij)"
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Applied to: Z_pre (softmax output from GCN)"
    p.font.size = Pt(12)
    p.font.italic = True
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Key Insight: Dual self-supervision"
    p.font.bold = True
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Both modules guided by same target distribution P"
    p.font.size = Pt(13)
    p.font.color.rgb = highlight_color
    
    # Slide 11: Training Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Training Strategy"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Two-Phase Training Process:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Phase 1: Pre-training (100 epochs)"
    p.font.bold = True
    p.font.size = Pt(16)
    
    phase1 = [
        "• Train ZINB autoencoder only",
        "• Loss: L_ZINB only",
        "• Learning rate: 0.001 (Adam optimizer)",
        "• Purpose: Good initialization for joint training"
    ]
    
    for item in phase1:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.space_after = Pt(3)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Phase 2: Joint Optimization (200 epochs)"
    p.font.bold = True
    p.font.size = Pt(16)
    
    phase2 = [
        "• Train all modules simultaneously",
        "• Loss: L_total (all 5 components)",
        "• Learning rate: 0.001 (Adam optimizer)",
        "• Initialize cluster centers with k-means",
        "• Update network parameters and cluster centers",
        "• Convergence: Stop when assignment changes < 0.1%"
    ]
    
    for item in phase2:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.space_after = Pt(3)
    
    # Slide 12: Evaluation Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Evaluation Techniques"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Primary Metrics:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "1. Normalized Mutual Information (NMI)"
    p.font.bold = True
    p.font.size = Pt(15)
    
    nmi_items = [
        "• Range: [0, 1], higher is better",
        "• Measures information shared between clusters and true labels",
        "• Formula: NMI = 2·I(C,K) / (H(C) + H(K))",
        "• Information-theoretic measure"
    ]
    
    for item in nmi_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.space_after = Pt(2)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "2. Adjusted Rand Index (ARI)"
    p.font.bold = True
    p.font.size = Pt(15)
    
    ari_items = [
        "• Range: [-1, 1], higher is better",
        "• Measures similarity adjusted for chance",
        "• Robust to cluster size imbalances",
        "• Widely accepted benchmark"
    ]
    
    for item in ari_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.space_after = Pt(2)
    
    # Slide 13: Experimental Setup
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Experimental Setup & Validation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Datasets (6 scRNA-seq datasets):"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    datasets = [
        "• QS Limb Muscle: 1,090 cells, 6 types (Smart-seq2)",
        "• Adam: 3,660 cells, 8 types (Drop-seq, kidney)",
        "• QS Diaphragm: 870 cells, 5 types (Smart-seq2)",
        "• QS Trachea: 1,350 cells, 4 types (Smart-seq2)",
        "• Romanov: 2,881 cells, 7 types (SMARTer, hypothalamus)",
        "• QX Limb Muscle: 3,909 cells, 6 types (10x)"
    ]
    
    for item in datasets:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.space_after = Pt(3)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Validation Strategy:"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    validation = [
        "• 5 repetitions per experiment",
        "• Tested with 500, 1000, 1500, 2000 genes",
        "• Comparison with 7 baseline methods",
        "• Ablation studies (w/o ZINB, Attention, Graph)",
        "• Parameter sensitivity analysis (α, k, HVG)"
    ]
    
    for item in validation:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.space_after = Pt(3)
    
    # Slide 14: Performance Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Performance Results"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Best Performance (QX Limb Muscle):"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = "• NMI: 0.9665 (vs 0.9556 best baseline)"
    p.font.size = Pt(14)
    p.font.color.rgb = highlight_color
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• ARI: 0.9829 (vs 0.9702 best baseline)"
    p.font.size = Pt(14)
    p.font.color.rgb = highlight_color
    p.font.bold = True
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Overall Performance:"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    overall = [
        "✓ Outperforms all 7 baseline methods on all 6 datasets",
        "✓ Average NMI: 0.8695 (vs 0.8392 for scDSC)",
        "✓ Average ARI: 0.8897 (vs 0.8497 for scDSC)",
        "✓ Largest improvement on QS Trachea (+10.2% NMI)",
        "✓ Consistent across multiple platforms",
        "✓ Robust to parameter variations"
    ]
    
    for item in overall:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.space_after = Pt(5)
    
    # Slide 15: Ablation Study
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Ablation Study Results"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Component Importance (Average across 6 datasets):"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    p.space_after = Pt(10)
    
    # Create table-like structure
    ablation_data = [
        ("scASDC (full model)", "0.8450", "0.8281", "Baseline"),
        ("w/o ZINB Loss", "0.8052", "0.7841", "-4.7% NMI, -5.3% ARI"),
        ("w/o Attention", "0.8343", "0.8215", "-1.3% NMI, -0.8% ARI"),
        ("w/o Graph Loss", "0.8103", "0.7609", "-4.1% NMI, -8.1% ARI")
    ]
    
    p = tf.add_paragraph()
    p.text = "Model Variant          NMI      ARI     Impact"
    p.font.size = Pt(12)
    p.font.name = "Courier New"
    p.font.bold = True
    p.space_after = Pt(5)
    
    for variant, nmi, ari, impact in ablation_data:
        p = tf.add_paragraph()
        text = f"{variant:22} {nmi:7} {ari:7} {impact}"
        p.text = text
        p.font.size = Pt(11)
        p.font.name = "Courier New"
        p.font.color.rgb = highlight_color if "full" in variant else RGBColor(0, 0, 0)
        p.space_after = Pt(3)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Key Findings:"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    findings = [
        "• All components are essential for optimal performance",
        "• ZINB Loss: Critical for distribution modeling",
        "• Graph Loss: Most important for ARI (-8.1%)",
        "• Attention: Consistent positive impact"
    ]
    
    for item in findings:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.space_after = Pt(3)
    
    # Slide 16: Baseline Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Comparison with Baseline Methods"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Seven Baseline Methods Compared:"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    p.space_after = Pt(8)
    
    baselines = [
        "1. DESC - Deep embedding with batch effect removal",
        "2. SDCN - Structural deep clustering network",
        "3. scDeepCluster - ZINB-based deep clustering",
        "4. DCA - Deep count autoencoder",
        "5. scDSC - Deep structural clustering (ZINB+GNN)",
        "6. DEC - Deep embedding clustering",
        "7. AttentionAE-sc - Attention-based autoencoder"
    ]
    
    for item in baselines:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.space_after = Pt(3)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Why scASDC is Superior:"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    superior = [
        "✓ Only method combining ZINB + GCN + Multi-head Attention",
        "✓ Layer-wise fusion (not simple concatenation)",
        "✓ Dual self-supervision on both modules",
        "✓ 100% win rate across all datasets"
    ]
    
    for item in superior:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.bold = True
        p.space_after = Pt(5)
    
    # Slide 17: Biological Validation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Biological Validation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Downstream Analysis Confirms Biological Relevance:"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "1. Differential Gene Expression"
    p.font.bold = True
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Identifies highly expressed genes per cluster"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Expression heatmaps show clear patterns"
    p.font.size = Pt(12)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Pathway Enrichment Analysis"
    p.font.bold = True
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• T Cell cluster enriched in:"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "  - T cell activation pathways"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  - T cell receptor signaling"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• B Cell cluster enriched in:"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "  - B cell activation pathways"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  - B cell receptor signaling"
    p.font.size = Pt(11)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. Visualization (UMAP)"
    p.font.bold = True
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Clear separation of cell populations"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Better than all baseline methods"
    p.font.size = Pt(12)
    
    # Slide 18: Key Innovations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Innovations of scASDC"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    innovations = [
        ("1. Dual-Module Architecture", 
         "Combines ZINB autoencoder (content) + GCN (structure)"),
        
        ("2. Layer-by-Layer Attention Fusion",
         "Multi-head attention at each layer, not simple concatenation"),
        
        ("3. Dual Self-Supervised Learning",
         "Applies clustering loss to both ZINB and GCN outputs"),
        
        ("4. ZINB Distribution Modeling",
         "Explicitly handles zero-inflation and overdispersion"),
        
        ("5. Oversmoothing Mitigation",
         "Attention fusion prevents GCN feature loss"),
        
        ("6. End-to-End Training",
         "Unified framework for representation + clustering")
    ]
    
    for title_text, desc in innovations:
        p = tf.add_paragraph()
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = accent_color
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.space_after = Pt(8)
    
    # Slide 19: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "scASDC: A Comprehensive Solution for scRNA-seq Clustering"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    p.space_after = Pt(15)
    
    summary_points = [
        "✓ Four integrated modules: ZINB + GCN + Attention + Self-supervised",
        "✓ Five complementary loss functions working together",
        "✓ Layer-wise attention fusion for intelligent integration",
        "✓ Handles scRNA-seq challenges: sparsity, zero-inflation, noise",
        "✓ Superior performance: Best on all 6 datasets",
        "✓ Biologically validated: Meaningful clusters and pathways",
        "✓ Robust and generalizable across platforms",
        "✓ Interpretable through attention weights and gene analysis"
    ]
    
    for item in summary_points:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    # Slide 20: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Key Takeaways:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "1. scASDC advances single-cell clustering through intelligent fusion of content and structural information"
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Multi-head attention mechanism is crucial for effective integration"
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. All components (ZINB, GCN, Attention, Self-supervision) are necessary for optimal performance"
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "4. Empirically validated across diverse datasets and biologically interpretable"
    p.font.size = Pt(14)
    p.space_after = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Future Work: Batch correction, spatial transcriptomics, trajectory inference"
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    # Save presentation
    output_path = "/home/frankfurt/projects/bio-informatics/Research Papers/scASDC_Analysis_Presentation.pptx"
    prs.save(output_path)
    print(f"✓ PowerPoint presentation created successfully!")
    print(f"✓ Saved to: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    try:
        output_file = create_scASDC_presentation()
        print("\n" + "="*60)
        print("PRESENTATION CONTENTS:")
        print("="*60)
        print("1. Title Slide")
        print("2. Overview of scASDC Method")
        print("3-6. Four Modules (ZINB, GCN, Attention, Self-supervised)")
        print("7-10. Five Loss Functions (detailed explanations)")
        print("11. Training Strategy")
        print("12. Evaluation Techniques")
        print("13. Experimental Setup")
        print("14. Performance Results")
        print("15. Ablation Study")
        print("16. Baseline Comparison")
        print("17. Biological Validation")
        print("18. Key Innovations")
        print("19. Summary")
        print("20. Conclusion")
        print("="*60)
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
